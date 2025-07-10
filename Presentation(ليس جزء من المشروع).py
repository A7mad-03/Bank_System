# sdk_bank_presentation.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def create_sdk_bank_presentation():
    # Create a presentation object
    prs = Presentation()

    # Logo path
    logo_path = r"C:\Users\Ahmad\Desktop\VS code\Final Project\SDK BANK (logo).png"

    # Function to add logo to a slide (top-right corner)
    def add_logo(slide):
        left = Inches(8.0)   # Adjust based on slide width (max 10")
        top = Inches(0.0)
        height = Inches(0.8)
        slide.shapes.add_picture(logo_path, left, top, height=height)

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "SDK Bank System"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)  # Dark blue
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle.text = "A Complete Banking Solution\nCreated with Python"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_logo(slide)

    # Slide 2: Features Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Key Features"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    features = [
        "User Account Management (Create/Login)",
        "Password Recovery System",
        "Deposit/Withdraw Funds",
        "Balance Checking",
        "Click Transfer between accounts",
        "Bill Payments (Electricity, Water, Telecom)",
        "Currency Conversion",
        "Streamlit Web Interface"
    ]

    for feature in features:
        p = content.text_frame.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True

    add_logo(slide)

    # Slide 3: Account Management
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Account Management"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    account_details = [
        "Create accounts with username, password, phone",
        "Security question for password recovery",
        "Login via username or phone number",
        "Secure password storage",
        "Session management"
    ]

    for detail in account_details:
        p = content.text_frame.add_paragraph()
        p.text = detail
        p.level = 0
        p.font.size = Pt(18)

    add_logo(slide)

    # Slide 4: Banking Operations
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Banking Operations"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    operations = [
        "Deposit: Add funds to account",
        "Withdraw: Remove funds (with balance check)",
        "Check Balance: View current balance",
        "Click Transfer: Send money to other users by phone",
        "All transactions logged"
    ]

    for op in operations:
        p = content.text_frame.add_paragraph()
        p.text = op
        p.level = 0
        p.font.size = Pt(18)

    add_logo(slide)

    # Slide 5: EFAWATEER Bill Payments
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Bill Payment System"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    bills = [
        "Electricity: Calculate bill based on device usage",
        "Water: Pay based on water consumption",
        "Telecom: Pay for calls and data usage",
        "All payments check account balance first",
        "Multi-step payment process"
    ]

    for bill in bills:
        p = content.text_frame.add_paragraph()
        p.text = bill
        p.level = 0
        p.font.size = Pt(18)

    add_logo(slide)

    # Slide 6: Currency Converter
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Currency Converter"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    converter_info = [
        "Convert Jordanian Dinar (JOD) to 20+ currencies",
        "Real-time exchange rates",
        "Simple dropdown interface",
        "Supports major world currencies",
        "Clear display of conversion results"
    ]

    for info in converter_info:
        p = content.text_frame.add_paragraph()
        p.text = info
        p.level = 0
        p.font.size = Pt(18)

    add_logo(slide)

    # Slide 7: Technical Implementation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Technical Implementation"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    tech_details = [
        "Python backend with Streamlit frontend",
        "Object-oriented design",
        "Modular components (Bank, Click, EFAWATEER)",
        "Session state management",
        "Input validation and error handling"
    ]

    for detail in tech_details:
        p = content.text_frame.add_paragraph()
        p.text = detail
        p.level = 0
        p.font.size = Pt(18)

    add_logo(slide)

    # Slide 8: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

    conclusion = [
        "Comprehensive banking solution",
        "User-friendly interface",
        "Secure transactions",
        "Extensible architecture",
        "Demonstrates Python's capabilities for financial apps"
    ]

    for point in conclusion:
        p = content.text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.bold = True

    add_logo(slide)

    # Save the presentation
    prs.save("SDK_Bank_Presentation.pptx")
    print("Presentation created successfully: SDK_Bank_Presentation.pptx")

if __name__ == "__main__":
    create_sdk_bank_presentation()
